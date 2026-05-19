"""Assemble a new .pptx by copying selected slides from source decks."""

import io
import re
import zipfile
from copy import deepcopy
from datetime import datetime
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.opc.packuri import PackURI
from pptx.util import Inches

from app.catalog import get_deck_by_id
from app.config import settings


_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_EMBED_ATTR = f"{{{_R_NS}}}embed"
_ID_ATTR = f"{{{_R_NS}}}id"
_LINK_ATTR = f"{{{_R_NS}}}link"

TEMPLATE_DIR = Path("slideTemplate")

# Any solid fill whose average RGB is below this is considered "dark"
_DARK_THRESHOLD = 140
# Scheme color names that map to dark in both source and template themes
_DARK_SCHEME_COLORS = {"dk1", "dk2"}


def _hex_luminance(hex_color: str) -> int:
    """Return average RGB (0–255) for a 6-hex-digit color string."""
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return (r + g + b) // 3
    except (ValueError, IndexError):
        return 0


def _shape_has_light_fill(sp_elem) -> bool:
    """Return True if the shape has an explicit solid fill that is light-colored."""
    spPr = sp_elem.find(f"{{{_P_NS}}}spPr")
    if spPr is None:
        return False
    solidFill = spPr.find(f"{{{_A_NS}}}solidFill")
    if solidFill is None:
        return False
    srgb = solidFill.find(f"{{{_A_NS}}}srgbClr")
    if srgb is not None:
        return _hex_luminance(srgb.get("val", "")) >= _DARK_THRESHOLD
    return False


def _set_run_white(rpr: etree._Element) -> None:
    """Overwrite or add a solid white fill on a run-properties element.

    <a:solidFill> must appear before font elements (<a:latin>, <a:ea>, etc.)
    in the OOXML schema. We always insert at position 0 so PowerPoint respects it.
    """
    for fill in rpr.findall(f"{{{_A_NS}}}solidFill"):
        rpr.remove(fill)
    sf = etree.Element(f"{{{_A_NS}}}solidFill")
    etree.SubElement(sf, f"{{{_A_NS}}}srgbClr").set("val", "FFFFFF")
    rpr.insert(0, sf)


def _fix_text_colors(new_slide) -> None:
    """
    Ensure text is visible on the dark template background.

    The source master defaulted to white text; the template master defaults to
    black. Any run property element without an explicit color would inherit black
    and become invisible. This function:
      - Covers <a:rPr>, <a:endParaRPr>, and field <a:rPr> (not just <a:r> runs).
      - Adds explicit white where there is no solidFill.
      - Replaces explicit dark sRGB / dark scheme colors (dk1, dk2) with white.
      - Adds white <a:buClr> to bullet paragraphs that have no bullet color.
      - Skips shapes with a light-colored explicit fill (dark text on white is fine).
    """
    _BULLET_TAGS = {
        f"{{{_A_NS}}}{t}"
        for t in ("buSzClr", "buSzPct", "buSzPts", "buFont", "buChar", "buAutoNum", "buNone", "buBlip")
    }

    for sp in new_slide._element.iter(f"{{{_P_NS}}}sp"):
        if _shape_has_light_fill(sp):
            continue

        # Fix every run-property element: <a:rPr> (in runs AND fields) and <a:endParaRPr>
        for rpr in sp.iter(f"{{{_A_NS}}}rPr", f"{{{_A_NS}}}endParaRPr"):
            solidFill = rpr.find(f"{{{_A_NS}}}solidFill")
            if solidFill is None:
                _set_run_white(rpr)
            else:
                srgb = solidFill.find(f"{{{_A_NS}}}srgbClr")
                scheme = solidFill.find(f"{{{_A_NS}}}schemeClr")
                if srgb is not None and _hex_luminance(srgb.get("val", "")) < _DARK_THRESHOLD:
                    _set_run_white(rpr)
                elif scheme is not None and scheme.get("val") in _DARK_SCHEME_COLORS:
                    _set_run_white(rpr)

        # Fix bullet colors in paragraph properties
        for pPr in sp.iter(f"{{{_A_NS}}}pPr"):
            buClr = pPr.find(f"{{{_A_NS}}}buClr")
            has_bullet = any(child.tag in _BULLET_TAGS for child in pPr)

            if buClr is None and has_bullet:
                # Insert white buClr before the first bullet sizing/font/char element
                insert_idx = next(
                    (i for i, c in enumerate(pPr) if c.tag in _BULLET_TAGS),
                    len(pPr),
                )
                bc = etree.Element(f"{{{_A_NS}}}buClr")
                etree.SubElement(bc, f"{{{_A_NS}}}srgbClr").set("val", "FFFFFF")
                pPr.insert(insert_idx, bc)
            elif buClr is not None:
                srgb = buClr.find(f"{{{_A_NS}}}srgbClr")
                scheme = buClr.find(f"{{{_A_NS}}}schemeClr")
                need_white = (
                    (srgb is not None and _hex_luminance(srgb.get("val", "")) < _DARK_THRESHOLD)
                    or (scheme is not None and scheme.get("val") in _DARK_SCHEME_COLORS)
                )
                if need_white:
                    for child in list(buClr):
                        buClr.remove(child)
                    etree.SubElement(buClr, f"{{{_A_NS}}}srgbClr").set("val", "FFFFFF")


def _find_template() -> Path | None:
    """Return the first .pptx found in the slideTemplate directory."""
    for f in sorted(TEMPLATE_DIR.glob("*.pptx")):
        return f
    return None


def _clear_slides(prs: Presentation) -> None:
    """Remove all slides from a presentation, keeping the slide master intact."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(f"{{{_R_NS}}}id")
        sldIdLst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass


def _blank_layout(prs: Presentation):
    """Return the cleanest blank slide layout from the presentation."""
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _ensure_unique_partname(part, package) -> None:
    """Rename part in-place if its partname already exists in package.

    python-pptx uses the partname as a dict key, so a collision silently
    overwrites the earlier part while both objects still get serialised to
    the ZIP, producing duplicate entries that PowerPoint treats as corrupt.
    """
    taken = {str(p.partname) for p in package._parts.values()}
    pn = str(part.partname)
    if pn not in taken:
        return
    # Strip trailing digits to find a stable base, then increment
    m = re.match(r'^(.*?)(\d+)(\.[^./]+)$', pn)
    if m:
        prefix, ext = m.group(1), m.group(3)
    else:
        m2 = re.match(r'^(.*?)(\.[^./]+)$', pn)
        if not m2:
            return
        prefix, ext = m2.group(1), m2.group(2)
    i = 2
    while True:
        candidate = f"{prefix}{i}{ext}"
        if candidate not in taken:
            part._partname = PackURI(candidate)
            return
        i += 1


def _remap_rids(slide_part, old_to_new: dict[str, str]) -> None:
    """Walk the slide XML and replace old rIds with new ones in-place."""
    for elem in slide_part._element.iter():
        for attr in (_EMBED_ATTR, _ID_ATTR, _LINK_ATTR):
            val = elem.get(attr)
            if val and val in old_to_new:
                elem.set(attr, old_to_new[val])


def _copy_slide(source_prs: Presentation, slide_index: int, target_prs: Presentation) -> None:
    """Copy a single slide (with images) from source_prs into target_prs."""
    source_slide = source_prs.slides[slide_index]

    blank_layout = _blank_layout(target_prs)
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Replace the new slide's spTree with a deep copy of the source slide's spTree
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)
    for child in source_slide.shapes._spTree:
        sp_tree.append(deepcopy(child))

    # Copy relationships (images, embedded objects, hyperlinks) and build rId remap table.
    # Skip slide layout / slide master / notes — those are managed by the target presentation.
    _SKIP_RELTYPES = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
    }
    old_to_new: dict[str, str] = {}
    for rel in source_slide.part.rels.values():
        if rel.reltype in _SKIP_RELTYPES:
            continue
        try:
            if rel.is_external:
                new_rid = new_slide.part.relate_to(rel.target_url, rel.reltype, is_external=True)
            else:
                _ensure_unique_partname(rel.target_part, target_prs.part.package)
                new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            if new_rid != rel.rId:
                old_to_new[rel.rId] = new_rid
        except Exception:
            pass

    if old_to_new:
        _remap_rids(new_slide.part, old_to_new)

    _fix_text_colors(new_slide)

    # Copy slide notes if present
    if source_slide.has_notes_slide:
        try:
            notes_text = source_slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                new_slide.notes_slide.notes_text_frame.text = notes_text
        except Exception:
            pass


def _deduplicate_zip(path: Path) -> None:
    """Rewrite the PPTX archive, keeping only the first occurrence of each ZIP entry.

    When slides are copied from multiple source .pptx files, python-pptx can write
    duplicate ZIP entries (e.g. two ppt/media/image1.png from two different sources).
    PowerPoint treats duplicate entries as corruption and triggers the repair prompt.
    """
    data = path.read_bytes()
    out = io.BytesIO()
    seen: set[str] = set()
    with zipfile.ZipFile(io.BytesIO(data)) as zin, \
            zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            if info.filename not in seen:
                seen.add(info.filename)
                zout.writestr(info, zin.read(info.filename))
    path.write_bytes(out.getvalue())


def _sanitize_filename(title: str) -> str:
    title = re.sub(r"[^\w\s-]", "", title).strip()
    title = re.sub(r"\s+", "_", title)
    return title[:60] or "Workshop"


def assemble_deck(
    selections: list[dict],
    catalog: dict,
    slides_dir: Path,
    output_path: Path,
) -> Path:
    """
    Build a new .pptx from the given selections, using the Kong slide template
    for theming so the output inherits the correct dark background and branding.
    """
    template_path = _find_template()
    if template_path:
        target_prs = Presentation(str(template_path))
        _clear_slides(target_prs)
    else:
        target_prs = Presentation()
        target_prs.slide_width = Inches(10)
        target_prs.slide_height = Inches(5.625)

    output_path.parent.mkdir(parents=True, exist_ok=True)

    for selection in selections:
        deck_meta = get_deck_by_id(catalog, selection["deck_id"])
        if not deck_meta:
            print(f"  Warning: deck '{selection['deck_id']}' not found in catalog, skipping.")
            continue

        pptx_path = slides_dir.parent / deck_meta["path"]
        if not pptx_path.exists():
            print(f"  Warning: file not found: {pptx_path}, skipping.")
            continue

        source_prs = Presentation(str(pptx_path))
        total_slides = len(source_prs.slides)

        for idx in selection.get("slide_indices", []):
            if 0 <= idx < total_slides:
                _copy_slide(source_prs, idx, target_prs)
            else:
                print(f"  Warning: slide index {idx} out of range for '{deck_meta['filename']}' ({total_slides} slides), skipping.")

    target_prs.save(str(output_path))
    _deduplicate_zip(output_path)
    return output_path


def build_output_path(title: str, output_dir: Path = settings.output_dir) -> Path:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_title = _sanitize_filename(title)
    return output_dir / f"{timestamp}_{safe_title}.pptx"
