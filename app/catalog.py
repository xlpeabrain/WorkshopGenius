"""Slide catalog: indexes the slides/ directory into a structured JSON catalog."""

import json
import os
import re
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from app.config import settings


def _slugify(text: str) -> str:
    text = text.lower().strip()
    text = re.sub(r"[^\w\s-]", "", text)
    text = re.sub(r"[\s_]+", "_", text)
    text = re.sub(r"-+", "_", text)
    return text[:60]


def _extract_slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    parts.append(line)
    return " ".join(parts)


def _get_slide_title(slide) -> str:
    """Return the title placeholder text, falling back to first text found."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()
        if title:
            return title
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text[:80]
    return "(no title)"


def _find_pptx_files(slides_dir: Path) -> list[Path]:
    """Walk slides_dir and return all .pptx files, excluding Repaired variants."""
    found = []
    for root, _dirs, files in os.walk(slides_dir):
        for fname in sorted(files):
            if fname.endswith(".pptx") and "Repaired" not in fname and not fname.startswith("."):
                found.append(Path(root) / fname)
    return found


def _build_deck_entry(pptx_path: Path, slides_dir: Path, deck_index: int) -> dict:
    """Extract metadata from a single .pptx file without calling Claude."""
    prs = Presentation(str(pptx_path))
    relative = pptx_path.relative_to(slides_dir.parent)
    topic = pptx_path.parent.name

    # Build stable slug from topic + sequential index
    deck_id = f"{_slugify(topic)}_{deck_index:02d}"

    slides_meta = []
    for i, slide in enumerate(prs.slides):
        title = _get_slide_title(slide)
        raw_text = _extract_slide_text(slide)
        slides_meta.append({
            "index": i,
            "title": title,
            "text_snippet": raw_text[:300],
        })

    return {
        "id": deck_id,
        "topic": topic,
        "filename": pptx_path.name,
        "path": str(relative),
        "slide_count": len(prs.slides),
        "tags": [],        # populated later by Claude (or left empty for basic mode)
        "description": "", # populated later by Claude
        "slides": slides_meta,
    }


def build_catalog(slides_dir: Path) -> dict:
    """Scan slides_dir and build a fresh catalog dict."""
    pptx_files = _find_pptx_files(slides_dir)
    decks = []
    for i, pptx_path in enumerate(pptx_files):
        print(f"  Indexing: {pptx_path.name}")
        entry = _build_deck_entry(pptx_path, slides_dir, i + 1)
        decks.append(entry)

    return {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "decks": decks,
    }


def _catalog_is_stale(slides_dir: Path, catalog_path: Path) -> bool:
    """Return True if catalog.json is missing or older than any .pptx file."""
    if not catalog_path.exists():
        return True
    catalog_mtime = catalog_path.stat().st_mtime
    for pptx_path in _find_pptx_files(slides_dir):
        if pptx_path.stat().st_mtime > catalog_mtime:
            return True
    return False


def load_or_rebuild_catalog(
    slides_dir: Path = settings.slides_dir,
    catalog_path: Path = settings.catalog_path,
) -> dict:
    """Load catalog.json if fresh, otherwise rebuild and persist it."""
    if _catalog_is_stale(slides_dir, catalog_path):
        print("Building slide catalog...")
        catalog = build_catalog(slides_dir)
        catalog_path.write_text(json.dumps(catalog, indent=2))
        print(f"Catalog saved: {len(catalog['decks'])} decks indexed.")
    else:
        catalog = json.loads(catalog_path.read_text())
        print(f"Loaded catalog: {len(catalog['decks'])} decks.")
    return catalog


def get_deck_by_id(catalog: dict, deck_id: str) -> dict | None:
    for deck in catalog["decks"]:
        if deck["id"] == deck_id:
            return deck
    return None


def get_catalog_for_prompt(catalog: dict) -> str:
    """Compact text representation of the catalog for Claude's system prompt."""
    lines = []
    for deck in catalog["decks"]:
        lines.append(f"[DECK: {deck['id']}]")
        lines.append(f"Topic: {deck['topic']}")
        lines.append(f"File: {deck['filename']}")
        lines.append(f"Slides: {deck['slide_count']} total")
        if deck.get("description"):
            lines.append(f"Description: {deck['description']}")
        if deck.get("tags"):
            lines.append(f"Tags: {', '.join(deck['tags'])}")
        lines.append("Key slides (index: title):")
        for s in deck["slides"]:
            snippet = f"  [{s['index']}] {s['title']}"
            if s.get("text_snippet"):
                preview = s["text_snippet"][:120].replace("\n", " ")
                snippet += f" — {preview}"
            lines.append(snippet)
        lines.append("")
    return "\n".join(lines)
